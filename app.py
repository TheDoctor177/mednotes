import os
import re
import json
import csv
from pathlib import Path
from typing import List, Dict, Any, Tuple

import streamlit as st
import whisper
import fitz  # PyMuPDF
from PIL import Image
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import yt_dlp

# -----------------------------
# Config
# -----------------------------
load_dotenv(".env", override=True)
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

st.set_page_config(page_title="MedNotes", page_icon="🩺", layout="wide")

if not OPENAI_API_KEY:
    st.error("OPENAI_API_KEY was not found. Put it in C:\\MedNotes\\.env as OPENAI_API_KEY=...")
    st.stop()

client = OpenAI(api_key=OPENAI_API_KEY)

APP_DIR = Path.cwd()
OUTPUTS_DIR = APP_DIR / "outputs"
OUTPUTS_DIR.mkdir(exist_ok=True)

SUPPORTED_UPLOADS = ["mp3", "wav", "m4a", "mp4", "mov", "mkv", "pdf", "pptx"]

# -----------------------------
# Helpers
# -----------------------------

def slugify(text: str) -> str:
    text = re.sub(r"[^a-zA-Z0-9_-]+", "-", text.strip())
    text = re.sub(r"-+", "-", text).strip("-")
    return text[:80] or "lecture"


def strip_code_fences(text: str) -> str:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-zA-Z0-9_]*\n", "", text)
        text = re.sub(r"\n```$", "", text)
    return text.strip()


def format_notes(notes: Any) -> str:
    if isinstance(notes, str):
        return notes
    if isinstance(notes, dict):
        parts = []
        for section, content in notes.items():
            parts.append(section)
            if isinstance(content, list):
                parts.extend([f"- {item}" for item in content])
            else:
                parts.append(str(content))
            parts.append("")
        return "\n".join(parts).strip()
    return str(notes)


def save_flashcards_csv(flashcards: List[Dict[str, str]], path: Path) -> None:
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        for card in flashcards:
            front = card.get("front", "").strip()
            back = card.get("back", "").strip()
            if front and back:
                writer.writerow([front, back])


def ensure_run_dir(title_hint: str) -> Path:
    run_name = slugify(title_hint)
    run_dir = OUTPUTS_DIR / run_name
    suffix = 1
    while run_dir.exists():
        suffix += 1
        run_dir = OUTPUTS_DIR / f"{run_name}-{suffix}"
    run_dir.mkdir(parents=True, exist_ok=True)
    return run_dir


@st.cache_resource
def load_whisper_model(model_name: str = "base"):
    return whisper.load_model(model_name)


def transcribe_media(file_path: str, model_name: str = "base") -> str:
    model = load_whisper_model(model_name)
    result = model.transcribe(file_path)
    return result["text"]


def download_youtube_media(url: str, out_dir: Path) -> Tuple[str, str]:
    outtmpl = str(out_dir / "youtube_media.%(ext)s")
    ydl_opts = {
        "format": "bestaudio/best",
        "outtmpl": outtmpl,
        "noplaylist": True,
        "quiet": True,
        "no_warnings": True,
    }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        info = ydl.extract_info(url, download=True)
        downloaded_path = Path(ydl.prepare_filename(info))
        if not downloaded_path.exists():
            matches = list(out_dir.glob("youtube_media.*"))
            if not matches:
                raise FileNotFoundError("Could not find downloaded YouTube media file.")
            downloaded_path = matches[0]
        title = info.get("title") or "YouTube lecture"
    return str(downloaded_path), title


def extract_pdf_text_and_images_with_placeholders(pdf_path: str, out_dir: Path) -> Tuple[str, Dict[str, str]]:
    """
    Extract text from PDF and insert [FIGURE_X] placeholders near the page
    where images are found. Returns:
      source_text, image_map
    where image_map = {"FIGURE_1": "/path/to/file.png", ...}
    """
    doc = fitz.open(pdf_path)
    texts = []
    image_map: Dict[str, str] = {}
    seen_hashes = set()
    figure_counter = 1

    for page_num, page in enumerate(doc, start=1):
        page_parts = [f"\n--- PDF page {page_num} ---\n"]

        page_text = page.get_text("text").strip()
        if page_text:
            page_parts.append(page_text)

        page_image_tokens = []
        for img_index, img in enumerate(page.get_images(full=True), start=1):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image.get("ext", "png")
            image_hash = hash(image_bytes)
            if image_hash in seen_hashes:
                continue
            seen_hashes.add(image_hash)

            token = f"FIGURE_{figure_counter}"
            image_path = out_dir / f"{token.lower()}_page_{page_num}_img_{img_index}.{image_ext}"
            with open(image_path, "wb") as f:
                f.write(image_bytes)

            image_map[token] = str(image_path)
            page_image_tokens.append(f"[{token}]")
            figure_counter += 1

        if page_image_tokens:
            page_parts.append("\nFigures on this page:\n" + "\n".join(page_image_tokens))

        texts.append("\n".join(page_parts))

    doc.close()
    return "\n".join(texts).strip(), image_map


def extract_pptx_text_and_images_with_placeholders(pptx_path: str, out_dir: Path) -> Tuple[str, Dict[str, str]]:
    """
    Extract text from PPTX and insert [FIGURE_X] placeholders on each slide
    after relevant text.
    """
    prs = Presentation(pptx_path)
    slide_texts = []
    image_map: Dict[str, str] = {}
    figure_counter = 1

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_parts = [f"\n--- Slide {slide_index} ---\n"]

        text_chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                text_chunks.append(shape.text.strip())

        if text_chunks:
            slide_parts.append("\n".join(text_chunks))

        slide_image_tokens = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext or "png"
                token = f"FIGURE_{figure_counter}"
                image_path = out_dir / f"{token.lower()}_slide_{slide_index}_{shape.shape_id}.{ext}"
                with open(image_path, "wb") as f:
                    f.write(image.blob)
                image_map[token] = str(image_path)
                slide_image_tokens.append(f"[{token}]")
                figure_counter += 1

        if slide_image_tokens:
            slide_parts.append("\nFigures on this slide:\n" + "\n".join(slide_image_tokens))

        slide_texts.append("\n".join(slide_parts))

    return "\n".join(slide_texts).strip(), image_map


def build_combined_source_text(extracted_text: str, manual_notes: str) -> str:
    parts = []
    if extracted_text.strip():
        parts.append("=== EXTRACTED SOURCE MATERIAL ===\n" + extracted_text.strip())
    if manual_notes.strip():
        parts.append("=== USER'S OWN NOTES ===\n" + manual_notes.strip())
    return "\n\n".join(parts).strip()


def build_prompt(source_text: str, card_style: str, include_input_figures: bool, create_diagrams: bool) -> str:
    flashcard_rule = (
        "Use cloze-deletion style when possible for memorisable facts."
        if card_style == "Cloze"
        else "Use concise question-answer cards."
    )

    figure_instructions = []
    if include_input_figures:
        figure_instructions.append(
            "If the source contains placeholders like [FIGURE_1], [FIGURE_2], etc., preserve them exactly in the notes at the most relevant locations. Do not rename them. Do not remove them unless completely irrelevant. Do not invent extra figure tokens unless they already exist in the source."
        )
    if create_diagrams:
        figure_instructions.append(
            "Create 1-2 simple Graphviz DOT diagrams that summarise mechanisms or diagnostic flow. Return them as strings in a list under the key 'graphviz_diagrams'."
        )

    figure_text = "\n".join(figure_instructions) if figure_instructions else "Do not include figure suggestions or diagrams."

    return f"""
You are a medical study assistant.

Based on the source material below, do all relevant tasks.

Important rules:
- Combine the extracted source material and the user's own notes into one coherent output.
- Prefer the most medically useful and high-yield formulation.
- Do not discard useful details from the user's own notes.
- If placeholders like [FIGURE_1] appear in the source, keep them exactly as written in the final notes where they fit best.
- The final notes should be clean and readable.

1. Create structured medical study notes with these sections:
- Definition
- Symptoms
- Pathophysiology
- Diagnosis
- Treatment / Management (only if present in the source)
- Key exam points

2. Create 20 high-yield Anki flashcards.
{flashcard_rule}

3. Create a short high-yield summary in 5-10 bullet points.

4. {figure_text}

Return valid JSON only, with this structure:
{{
  "title": "...",
  "notes": {{
    "Definition": "...",
    "Symptoms": ["..."],
    "Pathophysiology": ["..."],
    "Diagnosis": ["..."],
    "Treatment / Management": ["..."],
    "Key exam points": ["..."]
  }},
  "high_yield_summary": ["..."],
  "flashcards": [
    {{"front": "...", "back": "..."}}
  ],
  "graphviz_diagrams": ["digraph G {{ A -> B }}"]
}}

Source material:
{source_text}
""".strip()


def generate_notes_and_cards(source_text: str, card_style: str, include_input_figures: bool, create_diagrams: bool) -> Dict[str, Any]:
    prompt = build_prompt(source_text, card_style, include_input_figures, create_diagrams)
    response = client.responses.create(
        model="gpt-4.1-mini",
        input=prompt,
    )
    output_text = strip_code_fences(response.output_text)
    return json.loads(output_text)


def pick_images_to_display(image_paths: List[str], max_images: int = 5) -> List[str]:
    valid = []
    for path in image_paths:
        try:
            with Image.open(path) as img:
                width, height = img.size
                if width >= 200 and height >= 200:
                    valid.append(path)
        except Exception:
            continue
    return valid[:max_images]


def render_notes_with_inline_figures(notes_text: str, image_map: Dict[str, str], image_width: int = 500) -> None:
    """
    Renders note text and displays images inline wherever [FIGURE_X] appears.
    """
    token_pattern = r"(\[FIGURE_\d+\])"
    parts = re.split(token_pattern, notes_text)

    for part in parts:
        if not part:
            continue

        match = re.fullmatch(r"\[(FIGURE_\d+)\]", part.strip())
        if match:
            token = match.group(1)
            image_path = image_map.get(token)
            if image_path and Path(image_path).exists():
                st.image(image_path, caption=token, width=image_width)
            else:
                st.caption(f"{token} (image not found)")
        else:
            st.markdown(part)


def process_uploaded_file(uploaded_file, run_dir: Path, whisper_model: str) -> Dict[str, Any]:
    suffix = Path(uploaded_file.name).suffix.lower()
    save_path = run_dir / uploaded_file.name
    with open(save_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

    transcript = ""
    source_text = ""
    title = Path(uploaded_file.name).stem
    image_map: Dict[str, str] = {}

    if suffix in [".mp3", ".wav", ".m4a", ".mp4", ".mov", ".mkv"]:
        transcript = transcribe_media(str(save_path), whisper_model)
        source_text = transcript
    elif suffix == ".pdf":
        source_text, image_map = extract_pdf_text_and_images_with_placeholders(str(save_path), run_dir)
    elif suffix == ".pptx":
        source_text, image_map = extract_pptx_text_and_images_with_placeholders(str(save_path), run_dir)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")

    return {
        "title": title,
        "transcript": transcript,
        "source_text": source_text,
        "image_map": image_map,
        "saved_path": str(save_path),
    }


def process_youtube_link(url: str, run_dir: Path, whisper_model: str) -> Dict[str, Any]:
    media_path, title = download_youtube_media(url, run_dir)
    transcript = transcribe_media(media_path, whisper_model)
    return {
        "title": title,
        "transcript": transcript,
        "source_text": transcript,
        "image_map": {},
        "saved_path": media_path,
    }


# -----------------------------
# UI
# -----------------------------
st.title("🩺 MedNotes MVP+")
st.write(
    "Upload audio, video, PDF, or PPTX, or paste a YouTube link. You can also paste your own notes. Generate notes, Anki flashcards, and optional figures/diagrams."
)

with st.sidebar:
    st.header("Settings")
    whisper_model = st.selectbox("Whisper model", ["tiny", "base", "small"], index=1)
    card_style = st.selectbox("Flashcard style", ["QA", "Cloze"], index=0)
    figure_mode = st.selectbox(
        "Figures in output",
        ["None", "Reuse extracted input figures", "Create AI diagrams", "Both"],
        index=1,
    )
    show_transcript = st.checkbox("Show transcript/source text", value=True)
    max_images = st.slider("Max extracted figures to display", min_value=1, max_value=8, value=4)
    inline_image_width = st.slider("Inline figure width", min_value=250, max_value=900, value=500)

include_input_figures = figure_mode in ["Reuse extracted input figures", "Both"]
create_diagrams = figure_mode in ["Create AI diagrams", "Both"]

youtube_url = st.text_input("YouTube lecture link (optional)")
uploaded_file = st.file_uploader(
    "Or upload a file",
    type=SUPPORTED_UPLOADS,
    help="Supported: mp3, wav, m4a, mp4, mov, mkv, pdf, pptx",
)

manual_notes = st.text_area(
    "Optional: paste your own notes here",
    height=220,
    placeholder="Paste lecture notes, annotations, summaries, mnemonics, or bullet points here...",
)

if youtube_url and uploaded_file:
    st.warning("Use either a YouTube link or one uploaded file at a time for this MVP.")

if st.button("Generate MedNotes"):
    if not youtube_url and uploaded_file is None and not manual_notes.strip():
        st.error("Add either a YouTube link, an uploaded file, or your own notes.")
        st.stop()

    source_label = youtube_url if youtube_url else (uploaded_file.name if uploaded_file else "manual-notes")
    run_dir = ensure_run_dir(source_label)

    try:
        processed = {
            "title": "Manual notes",
            "transcript": "",
            "source_text": "",
            "image_map": {},
            "saved_path": "",
        }

        if youtube_url or uploaded_file is not None:
            with st.spinner("Extracting source material..."):
                if youtube_url:
                    processed = process_youtube_link(youtube_url, run_dir, whisper_model)
                else:
                    processed = process_uploaded_file(uploaded_file, run_dir, whisper_model)

        extracted_text = processed["source_text"].strip()
        transcript = processed["transcript"].strip()
        image_map = processed["image_map"]
        title = processed["title"]

        combined_source_text = build_combined_source_text(
            extracted_text=extracted_text,
            manual_notes=manual_notes,
        )

        if not combined_source_text:
            st.error("No usable text could be extracted or provided.")
            st.stop()

        with st.spinner("Generating notes, summary, flashcards, and optional diagrams..."):
            data = generate_notes_and_cards(
                source_text=combined_source_text,
                card_style=card_style,
                include_input_figures=include_input_figures,
                create_diagrams=create_diagrams,
            )

        notes = data.get("notes", {})
        notes_text = format_notes(notes)
        flashcards = data.get("flashcards", [])
        high_yield = data.get("high_yield_summary", [])
        graphviz_diagrams = data.get("graphviz_diagrams", []) if create_diagrams else []

        # Save outputs
        transcript_path = run_dir / "transcript_or_source.txt"
        notes_path = run_dir / "notes.txt"
        flashcards_path = run_dir / "flashcards.csv"
        json_path = run_dir / "full_output.json"
        image_map_path = run_dir / "image_map.json"
        manual_notes_path = run_dir / "manual_notes.txt"

        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(combined_source_text)
        with open(notes_path, "w", encoding="utf-8") as f:
            f.write(notes_text)
        with open(manual_notes_path, "w", encoding="utf-8") as f:
            f.write(manual_notes)
        save_flashcards_csv(flashcards, flashcards_path)
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        with open(image_map_path, "w", encoding="utf-8") as f:
            json.dump(image_map, f, ensure_ascii=False, indent=2)

        st.success("Done!")
        st.subheader(data.get("title") or title)

        tab1, tab2, tab3, tab4 = st.tabs(["Notes", "Flashcards", "Figures", "Raw text"])

        with tab1:
            st.markdown("### High-yield summary")
            if high_yield:
                for item in high_yield:
                    st.markdown(f"- {item}")

            st.markdown("### Structured study notes")
            render_notes_with_inline_figures(
                notes_text=notes_text,
                image_map=image_map,
                image_width=inline_image_width,
            )

            st.markdown("### Plain text version")
            st.text_area("Notes text", notes_text, height=300)

            with open(notes_path, "rb") as f:
                st.download_button("Download notes.txt", f, file_name="notes.txt")

        with tab2:
            st.markdown("### Flashcards preview")
            if not flashcards:
                st.info("No flashcards returned.")
            else:
                for i, card in enumerate(flashcards[:20], start=1):
                    st.markdown(f"**{i}. Front:** {card.get('front', '')}")
                    st.markdown(f"**Back:** {card.get('back', '')}")
                    st.markdown("---")

            with open(flashcards_path, "rb") as f:
                st.download_button("Download Anki CSV", f, file_name="flashcards.csv", mime="text/csv")

        with tab3:
            if include_input_figures:
                st.markdown("### Extracted input figures")
                all_images = list(image_map.values())
                picked = pick_images_to_display(all_images, max_images=max_images)
                if picked:
                    st.image(picked, caption=[Path(p).name for p in picked], width=350)
                else:
                    st.info("No reusable figures were extracted from the input.")

            if create_diagrams:
                st.markdown("### AI-generated diagrams")
                if graphviz_diagrams:
                    for idx, dot in enumerate(graphviz_diagrams, start=1):
                        st.markdown(f"**Diagram {idx}**")
                        try:
                            st.graphviz_chart(dot)
                        except Exception:
                            st.code(dot, language="dot")
                else:
                    st.info("No diagrams were returned.")

            if not include_input_figures and not create_diagrams:
                st.info("Figure mode is set to None.")

        with tab4:
            if show_transcript:
                st.text_area("Combined source text", combined_source_text, height=350)

            with open(transcript_path, "rb") as f:
                st.download_button("Download source text", f, file_name="transcript_or_source.txt")

        st.caption(f"Saved run folder: {run_dir}")

    except Exception as e:
        st.exception(e)

st.markdown("---")
st.markdown(
    "**Recommended packages**: streamlit, openai, python-dotenv, openai-whisper, yt-dlp, pymupdf, python-pptx, pillow"
)
